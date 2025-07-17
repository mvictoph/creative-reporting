import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import io
from PIL import Image
import warnings
from slack_sdk import WebClient
from slack_sdk.errors import SlackApiError
import cv2
import tempfile
import numpy as np

# Ignorer les avertissements OpenPyXL
warnings.filterwarnings('ignore', category=UserWarning, module='openpyxl')

# Configuration de la page
st.set_page_config(page_title="Amazon Creative Reporting", page_icon="📊", layout="wide")

# Définir les colonnes requises pour chaque type de rapport
CTR_REQUIRED_COLUMNS = ['Variant', 'Click-throughs', 'Impressions']
VCR_REQUIRED_COLUMNS = ['Variant', 'Video start', 'Video complete']

# Initialisation du client Slack
slack_token = os.environ.get('SLACK_TOKEN')
if slack_token:
    client = WebClient(token=slack_token)
else:
    st.warning("Slack token not found. Slack integration will not work.")
    client = None

def pixels_to_inches(pixels):
    return Inches(pixels / 96)

def resize_image(img, max_width, max_height):
    width, height = img.size
    aspect_ratio = width / height
    if width > max_width or height > max_height:
        if aspect_ratio > 1:
            new_width = max_width
            new_height = int(new_width / aspect_ratio)
        else:
            new_height = max_height
            new_width = int(new_height * aspect_ratio)
        return new_width, new_height
    return width, height

def validate_columns(df, report_type):
    required_columns = CTR_REQUIRED_COLUMNS if report_type == "CTR Report" else VCR_REQUIRED_COLUMNS
    missing_columns = [col for col in required_columns if col not in df.columns]
    return len(missing_columns) == 0, missing_columns

def extract_frame_from_video(video_file, time_in_seconds=3):
    # Créer un fichier temporaire pour la vidéo
    with tempfile.NamedTemporaryFile(suffix='.mp4', delete=False) as tmp_file:
        tmp_file.write(video_file.read())
        video_path = tmp_file.name

    try:
        # Ouvrir la vidéo
        cap = cv2.VideoCapture(video_path)
        
        # Aller à la frame souhaitée (3 secondes)
        cap.set(cv2.CAP_PROP_POS_MSEC, time_in_seconds * 1000)
        
        # Lire la frame
        success, frame = cap.read()
        
        if success:
            # Convertir BGR en RGB
            frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            # Convertir en image PIL
            image = Image.fromarray(frame_rgb)
            return image
        else:
            return None
    finally:
        cap.release()
        os.unlink(video_path)  # Supprimer le fichier temporaire

def calculate_metrics(df, variant, report_type):
    variant_data = df[df['Variant'] == variant]
    if report_type == "CTR Report":
        total_clicks = variant_data['Click-throughs'].sum()
        total_impressions = variant_data['Impressions'].sum()
        rate = total_clicks / total_impressions if total_impressions > 0 else 0
        return total_clicks, total_impressions, rate
    else:  # VCR Report
        video_starts = variant_data['Video start'].sum()
        video_completes = variant_data['Video complete'].sum()
        vcr = video_completes / video_starts if video_starts > 0 else 0
        return video_starts, video_completes, vcr

def apply_amazon_style(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Amazon Ember"
            run.font.color.rgb = RGBColor(0, 0, 0)

def find_best_matching_image(variant_name, image_files, similarity_threshold=0.6):
    def clean_string(s):
        return ''.join(c.lower() for c in s if c.isalnum())
    
    clean_variant = clean_string(variant_name)
    best_match = None
    highest_similarity = 0
    
    for img_name in image_files:
        clean_img_name = clean_string(os.path.splitext(img_name)[0])
        
        if clean_img_name in clean_variant or clean_variant in clean_img_name:
            similarity = 0.8
        else:
            variant_words = set(clean_variant.split())
            img_words = set(clean_img_name.split())
            common_words = variant_words.intersection(img_words)
            
            if common_words:
                similarity = len(common_words) / max(len(variant_words), len(img_words))
            else:
                common_chars = set(clean_img_name) & set(clean_variant)
                similarity = len(common_chars) / len(set(clean_img_name + clean_variant))
        
        if similarity > highest_similarity:
            highest_similarity = similarity
            best_match = img_name
    
    return (best_match, highest_similarity) if highest_similarity >= similarity_threshold else (None, 0)

def create_ppt_from_data(df, images_dict, report_type):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Configuration adaptée selon le type de rapport
    if report_type == "CTR Report":
        title_top = Inches(0.2)
        top_margin = Inches(1.0)
        spacing_y = Inches(2.8)  # Espacement vertical pour CTR
        spacing_x = Inches(3) + Inches(2)  # Espacement horizontal pour CTR
    else:  # VCR Report
        title_top = Inches(0.5)
        top_margin = Inches(1.8)  # Augmenté pour plus d'espace après le titre
        spacing_y = Inches(2.8)  # Espacement vertical ajusté pour VCR
        spacing_x = Inches(3) + Inches(1)  # Espacement horizontal réduit pour VCR
    
    # Titre du rapport
    title_box = slide.shapes.add_textbox(Inches(1), title_top, Inches(8), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Creatives Static Performance Report" if report_type == "CTR Report" else "Creatives Video Performance Report"
    title_frame.paragraphs[0].font.size = Pt(24)
    apply_amazon_style(title_frame)
    
    # Configuration de la grille
    items_per_row = 2
    left_margin = Inches(1)
    max_image_width = Inches(3)  # ~220 pixels
    max_image_height = Inches(2.5)  # ~180 pixels
    
    for index, variant in enumerate(df['Variant'].unique()):
        row_num = index // items_per_row
        col_num = index % items_per_row
        
        left = left_margin + (col_num * spacing_x)
        top = top_margin + (row_num * spacing_y)
        
        try:
            # Ajout de l'image
            matched_image, similarity = find_best_matching_image(variant, images_dict.keys())
            if matched_image:
                img = images_dict[matched_image]
                width, height = resize_image(img, 220, 180)
                image_width = pixels_to_inches(width)
                image_height = pixels_to_inches(height)
                
                img_byte_arr = io.BytesIO()
                img.save(img_byte_arr, format=img.format if img.format else 'JPEG')
                img_byte_arr.seek(0)
                
                if report_type == "VCR Report":
                    # Bordure noire simple sans ombre
                    border = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, 
                        left, 
                        top, 
                        image_width, 
                        image_height
                    )
                    border.line.color.rgb = RGBColor(0, 0, 0)
                    border.line.width = Pt(1.5)  # Épaisseur de la bordure
                    border.fill.background()  # Rend le rectangle transparent
                    border.shadow.inherit = False  # Désactive l'ombre
                
                slide.shapes.add_picture(img_byte_arr, left, top, width=image_width, height=image_height)
            
            # Ajout des informations textuelles
            text_top = top + image_height + Inches(0.1)
            
            # Texte Creative
            name_box = slide.shapes.add_textbox(left, text_top, image_width, Inches(0.2))
            name_frame = name_box.text_frame
            p = name_frame.paragraphs[0]
            p.text = f"Creative: {variant}"
            p.font.size = Pt(9)
            p.font.bold = True
            apply_amazon_style(name_frame)
            
            # Métriques adaptées selon le type de rapport
            val1, val2, rate = calculate_metrics(df, variant, report_type)
            metrics_box = slide.shapes.add_textbox(left, text_top + Inches(0.25), image_width, Inches(0.6))
            metrics_frame = metrics_box.text_frame
            
            if report_type == "CTR Report":
                metrics = [
                    f"Click-throughs: {val1:,}",
                    f"Impressions: {val2:,}",
                    f"CTR: {rate:.2%}"
                ]
            else:
                metrics = [
                    f"Video starts: {val1:,}",
                    f"Video completes: {val2:,}",
                    f"VCR: {rate:.2%}"
                ]
            
            for idx, metric in enumerate(metrics):
                if idx == 0:
                    p = metrics_frame.paragraphs[0]
                else:
                    p = metrics_frame.add_paragraph()
                p.text = metric
                p.font.size = Pt(9)
            
            apply_amazon_style(metrics_frame)
            
        except Exception as e:
            st.error(f"Error with variant {variant}: {str(e)}")
    
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer

def send_report_to_slack(file_buffer, filename, user_login):
    if not client:
        return False, "Slack client not initialized. Check your SLACK_TOKEN."
    
    try:
        # Personnaliser le message avec le login de l'utilisateur
        message = f"Here's the latest report from {user_login}!"

        # Upload file directly to channel
        response = client.files_upload_v2(
            channel="C095U79QZDL",
            file=file_buffer,
            filename=filename,
            initial_comment=message
        )
        
        if response['ok']:
            return True, f"Report successfully sent to Slack for {user_login}!"
        else:
            return False, f"Error in Slack response: {response}"
            
    except SlackApiError as e:
        return False, f"Error sending to Slack: {str(e)}"

def main():
    st.title("Creative Reporting Generator")
    st.markdown("---")
    
    # Sélection du type de rapport
    report_type = st.radio(
        "Select Report Type:",
        ("CTR Report", "VCR Report"),
        help="Choose CTR for static creatives or VCR for video creatives"
    )

    # Nouvelle section pour le nom du rapport
    report_name = st.text_input(
        "Creative Re Reporting Name",
        placeholder="e.g. Google_Creative_Report_Q1_2024",
        help="Choose a name for your report",
        key="report_name"
    )
    
    final_report_name = report_name if report_name else "Creative_Reporting"
    
    st.markdown("---")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader("📊 Excel File")
        excel_file = st.file_uploader("Upload your Excel file", type=['xlsx'])
        if ef excel_file:
            st.success("✅ Excel file successfully loaded")

    with col2:
        st.subheader("🖼️ Creative Files")
        if report_type == "CTR Report":
            allowed_types = ['jpg', 'jpeg', 'png']
            help_text = "Upload your creative images"
        else:  # VCR Report
            allowed_types = ['mp4']
            help_text = "Upload your video files"
            
        creative_files = st.file_uploader(
            help_text,
            type=allowed_types,
            accept_multiple_files=True
        )
        
        if creative_files:
            try:
                if report_type == "CTR Report":
                    images_dict = {
                        file.name: Image.open(file) 
                        for file in creative_files
                    }
                else:  # VCR Report
                    images_dict = {
                        file.name.replace('.mp4', '.jpg'): extract_frame_from_video(file)
                        for file in creative_files
                    }
                st.success(f"✅ {len(creative_files)} files loaded")
            except Exception as e:
                st.error(f"Error processing files: {str(e)}")

    if excel_file and creative_files:
        try:
            df = pd.read_excel(excel_file)
            
            # Valider les colonnes
            columns_valid, missing_columns = validate_columns(df, report_type)
            if not columns_valid:
                st.error(f"Missing required columns for {report_type}: {', '.join(missing_columns)}")
                return

            st.write(f"Number of variants detected: {len(df['Variant'].unique())}")
            
            col1, col2 = st.columns(2)
            
            with col1:
                if st.button("🚀 Generate PowerPoint Report"):
                    with st.spinner('Generating report...'):
                        pptx_buffer = create_ppt_from_data(df, images_dict, report_type)
                        st.success("Report generated successfully!")
                        
                        st.download_button(
                            label="📥 Download PowerPoint Report",
                            data=pptx_buffer,
                            file_name=f"{final_report_name}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )

            with col2:
                with st.expander("📤 Send to Slack (Optional)"):
                    user_login = st.text_input(
                        "Enter your Amazon login to send via Slack",
                        placeholder="e.g. jsmith",
                        help="Optional: Enter your login to send the report to Slack"
                    )
                    
                    if user_login:
                        if st.button("📤 Generate and Send to Slack"):
                            with st.spinner('Generating and sending to Slack...'):
                                pptx_buffer = create_ppt_from_data(df, images_dict, report_type)
                                success, message = send_report_to_slack(
                                    pptx_buffer,
                                    f"{final_report_name}.pptx",
                                    user_login
                                )
                                
                                if success:
                                    st.success(message)
                                else:
                                    st.error(message)
                
        except Exception as e:
            st.error(f"An error occurred: {str(e)}")

if __name__ == "__main__":
    main()
